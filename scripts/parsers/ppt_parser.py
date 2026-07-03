import logging
from pathlib import Path
from typing import List, Optional

# BaseParser와 ParsedContent는 기존 프로젝트 구조를 따른다고 가정합니다.
from parsers.base import BaseParser, ParsedContent

logger = logging.getLogger(__name__)

class PPTXParser(BaseParser):
    """python-pptx 라이브러리를 사용하여 PPTX 문서의 텍스트와 표를 추출하는 클래스입니다."""

    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    def parse(self, file_path: Path, **kwargs) -> ParsedContent:
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx 라이브러리가 설치되지 않았습니다. 'pip install python-pptx'를 실행하세요.")
            return ParsedContent(text="", parse_method="pptx_unavailable", confidence=0.0)

        file_path = Path(file_path)
        if not file_path.exists():
            return ParsedContent(text="", parse_method="file_not_found", confidence=0.0)

        try:
            # PPTX 파일 로드
            prs = Presentation(str(file_path))
        except Exception as e:
            logger.error(f"PPTX 파일을 여는 중 오류 발생: {file_path} -> {e}")
            return ParsedContent(text="", parse_method="open_failed", confidence=0.0)

        slide_texts: List[str] = []
        tables_markdown: List[str] = []

        # 슬라이드 순회
        for i, slide in enumerate(prs.slides, 1):
            current_slide_parts: List[str] = []

            # 슬라이드 내 모든 도형(Shape) 확인
            for shape in slide.shapes:
                # 1. 텍스트 프레임이 있는 경우 (텍스트 박스, 제목 등)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            current_slide_parts.append(text)

                # 2. 표(Table)가 있는 경우
                elif shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    if table_md:
                        tables_markdown.append(f"### 슬라이드 {i} 표\n{table_md}")

            # 슬라이드 단위로 텍스트 결합
            if current_slide_parts:
                combined_text = f"--- Slide {i} ---\n" + "\n".join(current_slide_parts)
                slide_texts.append(combined_text)

        # 결과 반환
        return ParsedContent(
            text="\n\n".join(slide_texts),
            tables=tables_markdown,
            metadata={
                "filename": file_path.name,
                "total_slides": len(prs.slides),
                "file_size": file_path.stat().st_size,
            },
            page_count=len(prs.slides),
            parse_method="python-pptx",
            confidence=1.0 if slide_texts else 0.5
        )

    def _table_to_markdown(self, table) -> str:
        """PPTX 표 객체를 마크다운 문자열로 변환합니다."""
        try:
            rows = []
            for i, row in enumerate(table.rows):
                # 셀 내부의 줄바꿈은 공백으로 치환하여 마크다운 깨짐 방지
                cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                
                # 첫 번째 행 이후 마크다운 구분선 추가
                if i == 0:
                    rows.append("|" + "|".join(["---" for _ in row.cells]) + "|")
            
            return "\n".join(rows)
        except Exception as e:
            logger.warning(f"표 변환 중 오류 발생: {e}")
            return ""